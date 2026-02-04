# cockatoo_v1/src/document_processing/extractors/pptx_extractor.py

"""
PowerPoint PPTX presentation extractor using python-pptx.
"""

import os
import logging
from typing import Dict, Any, List, Union, Optional
from pathlib import Path
import zipfile
import tempfile

from .base_extractor import BaseExtractor

logger = logging.getLogger(__name__)

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    logger.warning("python-pptx not installed. PPTX support will be limited.")


class PPTXExtractor(BaseExtractor):
    """
    Microsoft PowerPoint PPTX presentation extractor.
    """
    
    def __init__(self):
        """Initialize PPTX extractor."""
        super().__init__()
        if not HAS_PPTX:
            self.logger.warning(
                "python-pptx is not installed. Install with: pip install python-pptx"
            )
    
    def get_supported_formats(self) -> List[str]:
        """
        Return list of supported file formats.
        
        Returns:
            List of supported file extensions
        """
        return ['.pptx', '.ppt']
    
    def extract(self, file_path: Union[str, Path]) -> Dict[str, Any]:
        """
        Extract text and metadata from PowerPoint presentation.
        
        Args:
            file_path: Path to PPTX file
        
        Returns:
            Dictionary containing:
                - text: Extracted text content
                - metadata: Presentation metadata
                - slides: List of slides with content
                - notes: Speaker notes
                - layout_info: Slide layout information
                - media: Information about embedded media
        """
        file_path = Path(file_path)
        if not file_path.exists():
            raise FileNotFoundError(f"PowerPoint file not found: {file_path}")
        
        logger.info(f"Extracting PowerPoint: {file_path}")
        
        result = {
            "text": "",
            "metadata": self.get_basic_metadata(file_path),
            "slides": [],
            "notes": [],
            "layout_info": {},
            "media": [],
            "shapes": [],
            "charts": [],
            "tables": [],
            "extraction_method": "python-pptx" if HAS_PPTX else "zip_parsing",
        }
        
        if HAS_PPTX:
            try:
                result = self._extract_with_pptx(file_path, result)
            except Exception as e:
                logger.error(f"python-pptx extraction failed: {e}")
                if file_path.suffix.lower() == '.pptx':
                    result = self._extract_with_zip(file_path, result)
        else:
            if file_path.suffix.lower() == '.pptx':
                result = self._extract_with_zip(file_path, result)
            else:
                raise ImportError(
                    "python-pptx is required for .ppt files. "
                    "Install with: pip install python-pptx"
                )
        
        # Clean and combine all text
        all_text = []
        
        # Add slide texts
        for slide in result.get("slides", []):
            all_text.append(slide.get("text", ""))
        
        # Add notes
        for note in result.get("notes", []):
            all_text.append(note.get("text", ""))
        
        result["text"] = self.clean_text("\n\n".join(all_text))
        
        # Add language detection and summary
        if result["text"]:
            result["metadata"]["language"] = self.detect_language(result["text"])
            result["metadata"]["summary"] = self.extract_summary(result["text"])
        
        return result
    
    def _extract_with_pptx(self, file_path: Path, result: Dict[str, Any]) -> Dict[str, Any]:
        """
        Extract using python-pptx library.
        
        Args:
            file_path: Path to PPTX file
            result: Result dictionary to update
            
        Returns:
            Updated result dictionary
        """
        prs = Presentation(str(file_path))
        
        # Extract metadata
        core_properties = prs.core_properties
        result["metadata"].update({
            "title": core_properties.title or "",
            "subject": core_properties.subject or "",
            "author": core_properties.author or "",
            "comments": core_properties.comments or "",
            "keywords": core_properties.keywords or "",
            "category": core_properties.category or "",
            "created": str(core_properties.created) if core_properties.created else "",
            "modified": str(core_properties.modified) if core_properties.modified else "",
            "last_modified_by": core_properties.last_modified_by or "",
            "revision": core_properties.revision or "",
            "version": core_properties.version or "",
        })
        
        # Extract slides
        slides = []
        notes = []
        shapes = []
        charts = []
        tables = []
        
        for slide_idx, slide in enumerate(prs.slides):
            slide_text_parts = []
            slide_shapes = []
            
            # Extract from shapes
            for shape in slide.shapes:
                shape_info = {
                    "slide": slide_idx + 1,
                    "shape_id": shape.shape_id,
                    "shape_type": shape.shape_type,
                    "name": shape.name,
                    "text": "",
                    "has_text_frame": shape.has_text_frame,
                    "has_table": shape.has_table,
                    "has_chart": shape.has_chart,
                }
                
                # Extract text from shape
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    shape_text = ""
                    
                    for paragraph in text_frame.paragraphs:
                        paragraph_text = ""
                        for run in paragraph.runs:
                            paragraph_text += run.text
                        
                        shape_text += paragraph_text + "\n"
                    
                    shape_text = shape_text.strip()
                    shape_info["text"] = shape_text
                    
                    if shape_text:
                        slide_text_parts.append(shape_text)
                
                # Extract from tables
                if shape.has_table:
                    table_info = {
                        "slide": slide_idx + 1,
                        "shape_id": shape.shape_id,
                        "rows": len(shape.table.rows),
                        "columns": len(shape.table.columns),
                        "data": [],
                    }
                    
                    table_data = []
                    for row in shape.table.rows:
                        row_data = []
                        for cell in row.cells:
                            cell_text = cell.text_frame.text.strip() if cell.text_frame else ""
                            row_data.append(cell_text)
                        table_data.append(row_data)
                    
                    table_info["data"] = table_data
                    tables.append(table_info)
                    
                    # Add table text to slide text
                    table_text = "\n".join([" | ".join(row) for row in table_data])
                    slide_text_parts.append(table_text)
                
                # Extract chart information
                if shape.has_chart:
                    chart_info = {
                        "slide": slide_idx + 1,
                        "shape_id": shape.shape_id,
                        "chart_type": str(shape.chart.chart_type),
                    }
                    charts.append(chart_info)
                
                slide_shapes.append(shape_info)
            
            # Extract slide notes
            slide_notes = ""
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if notes_slide.notes_text_frame:
                    slide_notes = notes_slide.notes_text_frame.text
            
            # Combine slide text
            slide_text = "\n\n".join(slide_text_parts)
            
            slide_info = {
                "slide_number": slide_idx + 1,
                "slide_id": slide.slide_id,
                "slide_layout": slide.slide_layout.name,
                "text": slide_text,
                "shape_count": len(slide.shapes),
                "notes": slide_notes,
                "shapes": slide_shapes,
            }
            
            slides.append(slide_info)
            shapes.extend(slide_shapes)
            
            if slide_notes:
                notes.append({
                    "slide_number": slide_idx + 1,
                    "text": slide_notes,
                })
        
        result.update({
            "slides": slides,
            "notes": notes,
            "shapes": shapes,
            "charts": charts,
            "tables": tables,
        })
        
        # Extract layout information
        layout_info = {}
        for layout in prs.slide_layouts:
            layout_info[layout.name] = {
                "slide_count": len([s for s in slides if s["slide_layout"] == layout.name]),
                "layout_id": layout.slide_layout_id,
            }
        
        result["layout_info"] = layout_info
        
        # Count media elements
        media_count = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == 13:  # Picture shape type
                    media_count += 1
        
        result["media_count"] = media_count
        
        return result
    
    def _extract_with_zip(self, file_path: Path, result: Dict[str, Any]) -> Dict[str, Any]:
        """
        Fallback extraction by parsing PPTX as ZIP archive.
        
        Args:
            file_path: Path to PPTX file
            result: Result dictionary to update
            
        Returns:
            Updated result dictionary
        """
        slides = []
        
        try:
            with zipfile.ZipFile(file_path, 'r') as pptx_zip:
                # List slide files
                slide_files = [f for f in pptx_zip.namelist() 
                             if f.startswith('ppt/slides/slide') and f.endswith('.xml')]
                
                # Sort slides by number
                slide_files.sort(key=lambda x: int(x.split('slide')[1].split('.')[0]))
                
                for slide_file in slide_files:
                    with pptx_zip.open(slide_file) as f:
                        slide_xml = f.read().decode('utf-8', errors='ignore')
                        
                        # Extract text from XML
                        slide_text = self._extract_text_from_pptx_xml(slide_xml)
                        
                        # Extract slide number from filename
                        slide_num = slide_file.split('slide')[1].split('.')[0]
                        
                        slides.append({
                            "slide_number": int(slide_num),
                            "text": slide_text,
                            "file": slide_file,
                        })
                
                # Try to extract notes
                notes_files = [f for f in pptx_zip.namelist() 
                             if f.startswith('ppt/notesSlides/notesSlide') and f.endswith('.xml')]
                
                notes = []
                for notes_file in notes_files:
                    with pptx_zip.open(notes_file) as f:
                        notes_xml = f.read().decode('utf-8', errors='ignore')
                        notes_text = self._extract_text_from_pptx_xml(notes_xml)
                        
                        # Extract slide number from filename
                        slide_num = notes_file.split('notesSlide')[1].split('.')[0]
                        
                        notes.append({
                            "slide_number": int(slide_num),
                            "text": notes_text,
                        })
                
                result.update({
                    "slides": slides,
                    "notes": notes,
                })
                
        except zipfile.BadZipFile:
            logger.error(f"File is not a valid ZIP archive: {file_path}")
        except Exception as e:
            logger.error(f"Failed to parse PPTX as ZIP: {e}")
        
        return result
    
    def _extract_text_from_pptx_xml(self, xml_content: str) -> str:
        """
        Extract text from PPTX XML content.
        
        Args:
            xml_content: XML content as string
            
        Returns:
            Extracted text
        """
        import re
        
        # Remove namespaces
        xml_content = re.sub(r'xmlns[^=]*="[^"]*"', '', xml_content)
        
        # Extract text from <a:t> tags (text runs)
        text_pattern = r'<a:t[^>]*>(.*?)</a:t>'
        text_matches = re.findall(text_pattern, xml_content, re.DOTALL)
        
        # Clean text
        texts = []
        for text in text_matches:
            # Decode XML entities
            text = re.sub(r'&lt;', '<', text)
            text = re.sub(r'&gt;', '>', text)
            text = re.sub(r'&amp;', '&', text)
            text = re.sub(r'&quot;', '"', text)
            text = re.sub(r'&apos;', "'", text)
            
            # Remove other tags
            text = re.sub(r'<[^>]+>', '', text)
            
            if text.strip():
                texts.append(text.strip())
        
        return "\n".join(texts)
    
    def extract_slide_summaries(self, file_path: Union[str, Path]) -> List[Dict[str, Any]]:
        """
        Extract summaries for each slide.
        
        Args:
            file_path: Path to PPTX file
            
        Returns:
            List of slide summaries
        """
        result = self.extract(file_path)
        
        summaries = []
        for slide in result.get("slides", []):
            slide_text = slide.get("text", "")
            slide_notes = next((n["text"] for n in result.get("notes", []) 
                              if n["slide_number"] == slide["slide_number"]), "")
            
            combined_text = f"{slide_text}\n\n{slide_notes}".strip()
            
            if combined_text:
                summary = {
                    "slide_number": slide["slide_number"],
                    "title": self._extract_slide_title(combined_text),
                    "summary": self.extract_summary(combined_text, max_length=200),
                    "word_count": len(combined_text.split()),
                    "has_notes": bool(slide_notes),
                    "shape_count": slide.get("shape_count", 0),
                }
                summaries.append(summary)
        
        return summaries
    
    def _extract_slide_title(self, slide_text: str) -> str:
        """
        Extract title from slide text.
        
        Args:
            slide_text: Slide text content
            
        Returns:
            Extracted title
        """
        if not slide_text:
            return ""
        
        # Get first non-empty line
        lines = slide_text.strip().split('\n')
        for line in lines:
            line = line.strip()
            if line:
                # Check if line looks like a title (not too long, ends with punctuation)
                if len(line) < 100 and line.endswith(('.', '!', '?')):
                    return line
                return line
        
        return ""
    
    def get_presentation_stats(self, file_path: Union[str, Path]) -> Dict[str, Any]:
        """
        Get presentation statistics.
        
        Args:
            file_path: Path to PPTX file
            
        Returns:
            Presentation statistics
        """
        result = self.extract(file_path)
        
        slides = result.get("slides", [])
        notes = result.get("notes", [])
        
        stats = {
            "slide_count": len(slides),
            "total_words": sum(len(s.get("text", "").split()) for s in slides),
            "total_characters": sum(len(s.get("text", "")) for s in slides),
            "slides_with_notes": len(notes),
            "total_note_words": sum(len(n.get("text", "").split()) for n in notes),
            "average_words_per_slide": 0,
            "media_count": result.get("media_count", 0),
            "table_count": len(result.get("tables", [])),
            "chart_count": len(result.get("charts", [])),
            "shape_count": len(result.get("shapes", [])),
        }
        
        if slides:
            stats["average_words_per_slide"] = stats["total_words"] / len(slides)
        
        # Word frequency analysis (simplified)
        all_text = " ".join(s.get("text", "") for s in slides)
        words = all_text.lower().split()
        
        if words:
            from collections import Counter
            word_freq = Counter(words)
            stats["top_words"] = word_freq.most_common(10)
        
        return stats