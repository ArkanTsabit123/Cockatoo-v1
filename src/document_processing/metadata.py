# src/document_processing/metadata.py

"""Metadata extraction and management for document processing pipeline."""

import os
import re
import json
from pathlib import Path
from datetime import datetime
from typing import Dict, List, Optional, Any, Union
from dataclasses import dataclass, asdict, field
import logging
from enum import Enum

try:
    import pdfplumber
    HAS_PDFPLUMBER = True
except ImportError:
    HAS_PDFPLUMBER = False
    pdfplumber = None

try:
    from docx import Document as DocxDocument
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False
    DocxDocument = None

try:
    import ebooklib
    from ebooklib import epub
    HAS_EPUB = True
except ImportError:
    HAS_EPUB = False
    epub = None

try:
    import yaml
    HAS_YAML = True
except ImportError:
    HAS_YAML = False
    yaml = None

try:
    import frontmatter
    HAS_FRONTMATTER = True
except ImportError:
    HAS_FRONTMATTER = False
    frontmatter = None

logger = logging.getLogger(__name__)


class DocumentType(Enum):
    """Supported document types for metadata extraction."""
    PDF = "pdf"
    DOCX = "docx"
    TXT = "txt"
    EPUB = "epub"
    MD = "md"
    HTML = "html"
    JSON = "json"
    YAML = "yaml"
    PPTX = "pptx"
    XLSX = "xlsx"
    CSV = "csv"
    IMAGE = "image"
    UNKNOWN = "unknown"


@dataclass
class DocumentMetadata:
    """Comprehensive metadata structure for documents."""
    core_file_path: str
    core_file_name: str
    core_file_type: str
    core_file_size: int
    core_upload_date: Optional[datetime] = None
    
    content_title: Optional[str] = None
    content_author: Optional[str] = None
    content_date: Optional[datetime] = None
    content_publisher: Optional[str] = None
    content_language: Optional[str] = None
    content_abstract: Optional[str] = None
    content_keywords: str = ""
    content_categories: str = ""
    
    processing_status: str = "pending"
    processing_error: Optional[str] = None
    processing_chunk_count: int = 0
    processing_word_count: int = 0
    processing_page_count: Optional[int] = None
    
    technical_encoding: Optional[str] = None
    technical_created_date: Optional[datetime] = None
    technical_modified_date: Optional[datetime] = None
    technical_software: Optional[str] = None
    technical_version: Optional[str] = None
    
    custom_fields: Dict[str, Any] = field(default_factory=dict)
    
    tracking_is_indexed: bool = False
    tracking_indexed_at: Optional[datetime] = None
    tracking_last_accessed: Optional[datetime] = None
    tracking_access_count: int = 0
    
    def __post_init__(self):
        if self.core_upload_date is None:
            self.core_upload_date = datetime.now()
    
    def to_dict(self) -> Dict[str, Any]:
        result = asdict(self)
        
        for key, value in result.items():
            if isinstance(value, datetime):
                result[key] = value.isoformat()
        
        return result
    
    def to_json(self) -> str:
        return json.dumps(self.to_dict(), indent=2, ensure_ascii=False)
    
    @classmethod
    def from_dict(cls, data: Dict[str, Any]) -> 'DocumentMetadata':
        datetime_fields = ['core_upload_date', 'content_date', 'technical_created_date', 
                          'technical_modified_date', 'tracking_indexed_at', 'tracking_last_accessed']
        
        for field in datetime_fields:
            if field in data and data[field] and isinstance(data[field], str):
                try:
                    data[field] = datetime.fromisoformat(data[field].replace('Z', '+00:00'))
                except (ValueError, TypeError):
                    data[field] = None
        
        return cls(**data)


class MetadataExtractor:
    """Main class for extracting metadata from various document formats."""
    
    def __init__(self, config: Optional[Dict[str, Any]] = None):
        self.config = config or {}
        self.extractors = self._initialize_extractors()
        
        self.patterns = {
            'title': [
                r'^#{1,3}\s+(.+)$',
                r'<title>(.+?)</title>',
                r'^(?:[A-Z][A-Z\s]{10,})$',
                r'^[A-Z][a-z]+(?:\s+[A-Z][a-z]+){2,}$',
            ],
            'author': [
                r'by\s+([A-Z][a-z]+\s+[A-Z][a-z]+)',
                r'author[:\s]+([^\n]+)',
                r'©\s*\d{4}\s+([^\n]+)',
                r'written by\s+([^\n]+)',
            ],
            'date': [
                r'\b\d{4}-\d{2}-\d{2}\b',
                r'\b\d{2}/\d{2}/\d{4}\b',
                r'\b(?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)[a-z]* \d{1,2},? \d{4}\b',
                r'\b\d{4}\b',
            ],
            'email': [
                r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b'
            ],
            'url': [
                r'https?://(?:[-\w.]|(?:%[\da-fA-F]{2}))+'
            ],
            'isbn': [
                r'\b(?:ISBN(?:-1[03])?:? )?(?=[0-9X]{10}$|(?=(?:[0-9]+[- ]){3})[- 0-9X]{13}$|97[89][0-9]{10}$|(?=(?:[0-9]+[- ]){4})[- 0-9]{17}$)(?:97[89][- ]?)?[0-9]{1,5}[- ]?[0-9]+[- ]?[0-9]+[- ]?[0-9X]\b'
            ],
            'doi': [
                r'\b10\.\d{4,9}/[-._;()/:A-Z0-9]+\b'
            ]
        }
        
        self.language_patterns = {
            'en': [r'\bthe\b', r'\band\b', r'\bof\b'],
            'id': [r'\bdan\b', r'\byang\b', r'\bdi\b'],
            'es': [r'\bel\b', r'\bla\b', r'\by\b'],
            'fr': [r'\ble\b', r'\bla\b', r'\bet\b'],
            'de': [r'\bder\b', r'\bdie\b', r'\bden\b'],
        }
    
    def _initialize_extractors(self) -> Dict[DocumentType, Any]:
        extractors = {}
        
        if HAS_PDFPLUMBER:
            extractors[DocumentType.PDF] = self._extract_pdf_metadata
        else:
            logger.warning("pdfplumber not installed. PDF metadata extraction disabled.")
        
        if HAS_DOCX:
            extractors[DocumentType.DOCX] = self._extract_docx_metadata
        else:
            logger.warning("python-docx not installed. DOCX metadata extraction disabled.")
        
        if HAS_EPUB:
            extractors[DocumentType.EPUB] = self._extract_epub_metadata
        else:
            logger.warning("ebooklib not installed. EPUB metadata extraction disabled.")
        
        extractors[DocumentType.TXT] = self._extract_txt_metadata
        extractors[DocumentType.MD] = self._extract_md_metadata
        extractors[DocumentType.HTML] = self._extract_html_metadata
        extractors[DocumentType.JSON] = self._extract_json_metadata
        extractors[DocumentType.YAML] = self._extract_yaml_metadata
        extractors[DocumentType.CSV] = self._extract_csv_metadata
        extractors[DocumentType.PPTX] = self._extract_pptx_metadata
        extractors[DocumentType.XLSX] = self._extract_xlsx_metadata
        extractors[DocumentType.IMAGE] = self._extract_image_metadata
        
        return extractors
    
    def extract_metadata(self, file_path: Union[str, Path]) -> DocumentMetadata:
        file_path = Path(file_path)
        
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        file_stats = file_path.stat()
        file_type = self._detect_file_type(file_path)
        
        metadata = DocumentMetadata(
            core_file_path=str(file_path),
            core_file_name=file_path.name,
            core_file_type=file_type.value,
            core_file_size=file_stats.st_size,
            technical_created_date=datetime.fromtimestamp(file_stats.st_ctime),
            technical_modified_date=datetime.fromtimestamp(file_stats.st_mtime),
        )
        
        try:
            if file_type in self.extractors:
                extractor = self.extractors[file_type]
                format_metadata = extractor(file_path)
                
                for key, value in format_metadata.items():
                    if value is not None and value != "":
                        prefixed_key = f"content_{key}" if key in ['title', 'author', 'date', 'publisher', 
                                                                   'language', 'abstract', 'keywords', 'categories'] else key
                        if hasattr(metadata, prefixed_key):
                            setattr(metadata, prefixed_key, value)
                        else:
                            metadata.custom_fields[key] = value
            
            self._apply_fallback_extraction(file_path, metadata)
            
            if metadata.processing_word_count == 0:
                metadata.processing_word_count = self._count_words_in_file(file_path)
            
            if not metadata.content_language:
                metadata.content_language = self._detect_language(file_path)
            
            if not metadata.content_keywords:
                metadata.content_keywords = self._extract_keywords(file_path, metadata.content_title)
            
            metadata.processing_status = "completed"
            
        except Exception as e:
            metadata.processing_status = "failed"
            metadata.processing_error = str(e)
            logger.error(f"Error extracting metadata from {file_path}: {e}")
        
        return metadata
    
    def _detect_file_type(self, file_path: Path) -> DocumentType:
        extension = file_path.suffix.lower()
        
        type_mapping = {
            '.pdf': DocumentType.PDF,
            '.docx': DocumentType.DOCX,
            '.doc': DocumentType.DOCX,
            '.txt': DocumentType.TXT,
            '.text': DocumentType.TXT,
            '.epub': DocumentType.EPUB,
            '.md': DocumentType.MD,
            '.markdown': DocumentType.MD,
            '.html': DocumentType.HTML,
            '.htm': DocumentType.HTML,
            '.json': DocumentType.JSON,
            '.yaml': DocumentType.YAML,
            '.yml': DocumentType.YAML,
            '.pptx': DocumentType.PPTX,
            '.ppt': DocumentType.PPTX,
            '.xlsx': DocumentType.XLSX,
            '.xls': DocumentType.XLSX,
            '.csv': DocumentType.CSV,
            '.jpg': DocumentType.IMAGE,
            '.jpeg': DocumentType.IMAGE,
            '.png': DocumentType.IMAGE,
            '.gif': DocumentType.IMAGE,
            '.bmp': DocumentType.IMAGE,
            '.tiff': DocumentType.IMAGE,
            '.tif': DocumentType.IMAGE,
        }
        
        return type_mapping.get(extension, DocumentType.UNKNOWN)
    
    def _extract_pdf_metadata(self, file_path: Path) -> Dict[str, Any]:
        metadata = {}
        
        try:
            with pdfplumber.open(file_path) as pdf:
                if pdf.metadata:
                    pdf_info = pdf.metadata
                    
                    field_mapping = {
                        'Title': 'title',
                        'Author': 'author',
                        'Subject': 'abstract',
                        'Keywords': 'keywords',
                        'Producer': 'software',
                        'Creator': 'software',
                        'CreationDate': 'created_date',
                        'ModDate': 'modified_date',
                    }
                    
                    for pdf_field, our_field in field_mapping.items():
                        if pdf_field in pdf_info and pdf_info[pdf_field]:
                            value = pdf_info[pdf_field]
                            
                            if 'Date' in pdf_field and isinstance(value, str):
                                try:
                                    value = self._parse_pdf_date(value)
                                except (ValueError, TypeError):
                                    pass
                            
                            metadata[our_field] = value
                
                metadata['page_count'] = len(pdf.pages)
                
                if 'title' not in metadata or not metadata['title']:
                    if pdf.pages:
                        first_page = pdf.pages[0]
                        text = first_page.extract_text()
                        if text:
                            lines = text.split('\n')
                            if lines:
                                potential_title = lines[0].strip()
                                if len(potential_title) < 200 and self._looks_like_title(potential_title):
                                    metadata['title'] = potential_title
                
        except Exception as e:
            logger.warning(f"Error extracting PDF metadata from {file_path}: {e}")
        
        return metadata
    
    def _parse_pdf_date(self, pdf_date: str) -> Optional[datetime]:
        try:
            if pdf_date.startswith('D:'):
                date_str = pdf_date[2:]
                year = int(date_str[0:4])
                month = int(date_str[4:6]) if len(date_str) >= 6 else 1
                day = int(date_str[6:8]) if len(date_str) >= 8 else 1
                hour = int(date_str[8:10]) if len(date_str) >= 10 else 0
                minute = int(date_str[10:12]) if len(date_str) >= 12 else 0
                second = int(date_str[12:14]) if len(date_str) >= 14 else 0
                
                return datetime(year, month, day, hour, minute, second)
        except (ValueError, IndexError):
            pass
        
        return None
    
    def _extract_docx_metadata(self, file_path: Path) -> Dict[str, Any]:
        metadata = {}
        
        try:
            doc = DocxDocument(file_path)
            
            core_props = doc.core_properties
            
            if core_props.title:
                metadata['title'] = core_props.title
            if core_props.author:
                metadata['author'] = core_props.author
            if core_props.subject:
                metadata['abstract'] = core_props.subject
            if core_props.keywords:
                if isinstance(core_props.keywords, list):
                    metadata['keywords'] = ", ".join(core_props.keywords)
                else:
                    metadata['keywords'] = core_props.keywords
            if core_props.language:
                metadata['language'] = core_props.language
            if core_props.created:
                metadata['created_date'] = core_props.created
            if core_props.modified:
                metadata['modified_date'] = core_props.modified
            if core_props.last_modified_by:
                metadata['software'] = f"Microsoft Word ({core_props.last_modified_by})"
            
            metadata['page_count'] = len(doc.element.xpath('//w:pgSz'))
            
            if 'title' not in metadata or not metadata['title']:
                for paragraph in doc.paragraphs[:10]:
                    if paragraph.text and len(paragraph.text) < 200:
                        if paragraph.style.name.startswith('Heading') or self._looks_like_title(paragraph.text):
                            metadata['title'] = paragraph.text.strip()
                            break
            
        except Exception as e:
            logger.warning(f"Error extracting DOCX metadata from {file_path}: {e}")
        
        return metadata
    
    def _extract_epub_metadata(self, file_path: Path) -> Dict[str, Any]:
        metadata = {}
        
        try:
            book = epub.read_epub(file_path)
            
            dc_metadata = book.get_metadata('DC', {})
            
            dc_mapping = {
                'title': 'title',
                'creator': 'author',
                'description': 'abstract',
                'publisher': 'publisher',
                'date': 'date',
                'language': 'language',
                'identifier': 'isbn',
                'subject': 'keywords',
            }
            
            for dc_field, our_field in dc_mapping.items():
                if dc_field in dc_metadata:
                    values = dc_metadata[dc_field]
                    if values:
                        if our_field == 'keywords':
                            if isinstance(values, list):
                                metadata[our_field] = ", ".join(str(v) for v in values)
                            else:
                                metadata[our_field] = str(values)
                        else:
                            metadata[our_field] = values[0] if isinstance(values, list) else values
            
            metadata['page_count'] = len(book.spine)
            
        except Exception as e:
            logger.warning(f"Error extracting EPUB metadata from {file_path}: {e}")
        
        return metadata
    
    def _extract_txt_metadata(self, file_path: Path) -> Dict[str, Any]:
        metadata = {}
        
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read(5000)
                
                for field, patterns in self.patterns.items():
                    for pattern in patterns:
                        matches = re.findall(pattern, content, re.IGNORECASE | re.MULTILINE)
                        if matches:
                            if field == 'keywords':
                                if isinstance(matches, list):
                                    metadata[field] = ", ".join(str(m) for m in list(set(matches))[:10])
                                else:
                                    metadata[field] = str(matches)
                            elif field == 'categories':
                                if isinstance(matches, list):
                                    metadata[field] = ", ".join(str(m) for m in list(set(matches))[:5])
                                else:
                                    metadata[field] = str(matches)
                            else:
                                metadata[field] = matches[0]
                            break
                
                if 'title' not in metadata:
                    lines = content.split('\n')
                    for i, line in enumerate(lines[:10]):
                        line = line.strip()
                        if line and self._looks_like_title(line):
                            metadata['title'] = line
                            break
                
        except Exception as e:
            logger.warning(f"Error extracting TXT metadata from {file_path}: {e}")
        
        return metadata
    
    def _extract_md_metadata(self, file_path: Path) -> Dict[str, Any]:
        metadata = {}
        
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            if HAS_FRONTMATTER:
                try:
                    post = frontmatter.loads(content)
                    if post.metadata:
                        for key, value in post.metadata.items():
                            if key in ['title', 'author', 'date', 'categories', 'tags', 'keywords']:
                                if key in ['categories', 'tags', 'keywords']:
                                    if isinstance(value, list):
                                        metadata[key] = ", ".join(str(v) for v in value)
                                    else:
                                        metadata[key] = str(value)
                                else:
                                    metadata[key] = value
                except Exception:
                    pass
            
            if not metadata and content.startswith('---'):
                yaml_end = content.find('---', 3)
                if yaml_end != -1:
                    yaml_content = content[3:yaml_end]
                    if HAS_YAML:
                        try:
                            yaml_metadata = yaml.safe_load(yaml_content)
                            if yaml_metadata:
                                for key, value in yaml_metadata.items():
                                    if key in ['title', 'author', 'date', 'categories', 'tags', 'keywords']:
                                        if key in ['categories', 'tags', 'keywords']:
                                            if isinstance(value, list):
                                                metadata[key] = ", ".join(str(v) for v in value)
                                            else:
                                                metadata[key] = str(value)
                                        else:
                                            metadata[key] = value
                        except yaml.YAMLError:
                            pass
            
            if 'title' not in metadata:
                heading_match = re.search(r'^#\s+(.+)$', content, re.MULTILINE)
                if heading_match:
                    metadata['title'] = heading_match.group(1).strip()
            
        except Exception as e:
            logger.warning(f"Error extracting MD metadata from {file_path}: {e}")
        
        return metadata
    
    def _extract_html_metadata(self, file_path: Path) -> Dict[str, Any]:
        metadata = {}
        
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            
            meta_patterns = {
                'title': r'<title>(.+?)</title>',
                'description': r'<meta\s+name=["\']description["\']\s+content=["\'](.+?)["\']',
                'keywords': r'<meta\s+name=["\']keywords["\']\s+content=["\'](.+?)["\']',
                'author': r'<meta\s+name=["\']author["\']\s+content=["\'](.+?)["\']',
                'date': r'<meta\s+name=["\']date["\']\s+content=["\'](.+?)["\']',
                'og:title': r'<meta\s+property=["\']og:title["\']\s+content=["\'](.+?)["\']',
            }
            
            for field, pattern in meta_patterns.items():
                match = re.search(pattern, content, re.IGNORECASE)
                if match:
                    value = match.group(1).strip()
                    if field == 'keywords':
                        metadata['keywords'] = ", ".join(k.strip() for k in value.split(','))
                    elif field == 'og:title' and 'title' not in metadata:
                        metadata['title'] = value
                    elif field != 'og:title':
                        metadata[field] = value
            
            if 'title' not in metadata:
                h1_match = re.search(r'<h1[^>]*>(.+?)</h1>', content, re.IGNORECASE)
                if h1_match:
                    h1_text = re.sub(r'<[^>]+>', '', h1_match.group(1)).strip()
                    if h1_text:
                        metadata['title'] = h1_text
            
        except Exception as e:
            logger.warning(f"Error extracting HTML metadata from {file_path}: {e}")
        
        return metadata
    
    def _extract_json_metadata(self, file_path: Path) -> Dict[str, Any]:
        metadata = {}
        
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            
            common_fields = {
                'title': ['title', 'name', 'headline', 'subject'],
                'author': ['author', 'creator', 'by', 'writer'],
                'date': ['date', 'created', 'published', 'timestamp'],
                'abstract': ['abstract', 'summary', 'description', 'content'],
                'keywords': ['keywords', 'tags', 'categories', 'topics'],
            }
            
            def find_in_dict(data_dict, field_names):
                if isinstance(data_dict, dict):
                    for field in field_names:
                        if field in data_dict:
                            value = data_dict[field]
                            if value:
                                return value
                    for value in data_dict.values():
                        if isinstance(value, dict):
                            result = find_in_dict(value, field_names)
                            if result:
                                return result
                return None
            
            for metadata_field, json_fields in common_fields.items():
                value = find_in_dict(data, json_fields)
                if value:
                    if metadata_field == 'keywords' and isinstance(value, list):
                        metadata[metadata_field] = ", ".join(str(v) for v in value[:10])
                    else:
                        metadata[metadata_field] = value
            
        except (json.JSONDecodeError, UnicodeDecodeError) as e:
            logger.warning(f"Error parsing JSON metadata from {file_path}: {e}")
        except Exception as e:
            logger.warning(f"Error extracting JSON metadata from {file_path}: {e}")
        
        return metadata
    
    def _extract_yaml_metadata(self, file_path: Path) -> Dict[str, Any]:
        metadata = {}
        
        if not HAS_YAML:
            return metadata
        
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                data = yaml.safe_load(f)
            
            if isinstance(data, dict):
                common_fields = {
                    'title': ['title', 'name', 'Title', 'Name'],
                    'author': ['author', 'Author', 'by', 'By'],
                    'date': ['date', 'Date', 'created', 'Created'],
                    'description': ['description', 'Description', 'summary', 'Summary'],
                    'tags': ['tags', 'Tags', 'keywords', 'Keywords'],
                }
                
                for metadata_field, yaml_fields in common_fields.items():
                    for field in yaml_fields:
                        if field in data:
                            value = data[field]
                            if value:
                                if metadata_field == 'tags':
                                    if isinstance(value, list):
                                        metadata['keywords'] = ", ".join(str(v) for v in value)
                                    elif isinstance(value, str):
                                        metadata['keywords'] = value
                                else:
                                    metadata[metadata_field] = value
                                break
            
        except yaml.YAMLError as e:
            logger.warning(f"Error parsing YAML metadata from {file_path}: {e}")
        except Exception as e:
            logger.warning(f"Error extracting YAML metadata from {file_path}: {e}")
        
        return metadata
    
    def _extract_csv_metadata(self, file_path: Path) -> Dict[str, Any]:
        metadata = {}
        
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                lines = []
                for i in range(5):
                    line = f.readline()
                    if not line:
                        break
                    lines.append(line)
                
                if lines:
                    first_line = lines[0].strip()
                    if first_line:
                        columns = first_line.split(',')
                        if len(columns) > 1 and all(col.strip() for col in columns):
                            metadata['columns'] = ", ".join(col.strip() for col in columns)
                            metadata['row_count_estimate'] = sum(1 for _ in open(file_path)) - 1
        except Exception as e:
            logger.warning(f"Error extracting CSV metadata from {file_path}: {e}")
        
        return metadata
    
    def _extract_pptx_metadata(self, file_path: Path) -> Dict[str, Any]:
        metadata = {}
        
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            
            metadata['page_count'] = len(prs.slides)
            
            if prs.slides:
                first_slide = prs.slides[0]
                for shape in first_slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text = shape.text.strip()
                        if text and self._looks_like_title(text):
                            metadata['title'] = text
                            break
            
        except ImportError:
            logger.warning("python-pptx not installed. PPTX metadata extraction disabled.")
        except Exception as e:
            logger.warning(f"Error extracting PPTX metadata from {file_path}: {e}")
        
        return metadata
    
    def _extract_xlsx_metadata(self, file_path: Path) -> Dict[str, Any]:
        metadata = {}
        
        try:
            import pandas as pd
            
            excel_file = pd.ExcelFile(file_path)
            
            metadata['sheets'] = ", ".join(excel_file.sheet_names)
            
            try:
                df = pd.read_excel(file_path, sheet_name=0, nrows=5)
                metadata['columns'] = ", ".join(str(col) for col in df.columns.tolist())
                metadata['row_count_estimate'] = sum(1 for _ in open(file_path, 'rb')) // 100
            except:
                pass
            
        except ImportError:
            logger.warning("pandas not installed. Excel metadata extraction disabled.")
        except Exception as e:
            logger.warning(f"Error extracting Excel metadata from {file_path}: {e}")
        
        return metadata
    
    def _extract_image_metadata(self, file_path: Path) -> Dict[str, Any]:
        metadata = {}
        
        try:
            from PIL import Image
            import PIL.ExifTags
            
            with Image.open(file_path) as img:
                metadata['width'] = img.width
                metadata['height'] = img.height
                metadata['format'] = img.format
                metadata['mode'] = img.mode
                
                if hasattr(img, '_getexif') and img._getexif():
                    exif = img._getexif()
                    exif_tags = {
                        PIL.ExifTags.TAGS.get(tag, tag): value
                        for tag, value in exif.items()
                    }
                    
                    exif_mapping = {
                        'DateTime': 'created_date',
                        'Artist': 'author',
                        'ImageDescription': 'description',
                        'Copyright': 'copyright',
                    }
                    
                    for exif_tag, our_field in exif_mapping.items():
                        if exif_tag in exif_tags:
                            value = exif_tags[exif_tag]
                            if value:
                                metadata[our_field] = value
                
        except ImportError:
            logger.warning("Pillow not installed. Image metadata extraction disabled.")
        except Exception as e:
            logger.warning(f"Error extracting image metadata from {file_path}: {e}")
        
        return metadata
    
    def _apply_fallback_extraction(self, file_path: Path, metadata: DocumentMetadata):
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read(10000)
            
            if not metadata.content_title:
                metadata.content_title = self._extract_title_from_content(content)
            
            if not metadata.content_author:
                metadata.content_author = self._extract_author_from_content(content)
            
            if not metadata.content_date:
                metadata.content_date = self._extract_date_from_content(content)
            
            self._extract_additional_metadata(content, metadata)
            
        except Exception as e:
            logger.warning(f"Error in fallback extraction for {file_path}: {e}")
    
    def _extract_title_from_content(self, content: str) -> Optional[str]:
        lines = content.split('\n')
        
        for line in lines[:20]:
            line = line.strip()
            if line and self._looks_like_title(line):
                return line
        
        first_lines = lines[:50]
        if first_lines:
            non_empty_lines = [line.strip() for line in first_lines if line.strip()]
            if non_empty_lines:
                return max(non_empty_lines, key=len)
        
        return None
    
    def _looks_like_title(self, text: str) -> bool:
        if not text or len(text) > 200:
            return False
        
        words = text.split()
        
        if len(words) < 2:
            return False
        
        capitalized_words = sum(1 for word in words if word and word[0].isupper())
        capitalization_ratio = capitalized_words / len(words) if words else 0
        
        non_title_patterns = [
            r'^\d+\.',
            r'^[A-Za-z]\)',
            r'^\s*$',
            r'^[\-\*\+]\s',
        ]
        
        for pattern in non_title_patterns:
            if re.match(pattern, text):
                return False
        
        return 0.3 < capitalization_ratio < 0.9
    
    def _extract_author_from_content(self, content: str) -> Optional[str]:
        author_patterns = [
            r'by\s+([A-Z][a-z]+(?:\s+[A-Z][a-z]+){1,2})',
            r'Author[:\s]+([^\n]+)',
            r'©\s*\d{4}\s+([^\n]+)',
            r'Written by\s+([^\n]+)',
            r'Created by\s+([^\n]+)',
        ]
        
        for pattern in author_patterns:
            match = re.search(pattern, content, re.IGNORECASE)
            if match:
                author = match.group(1).strip()
                if len(author.split()) <= 4:
                    return author
        
        return None
    
    def _extract_date_from_content(self, content: str) -> Optional[datetime]:
        date_patterns = [
            r'\b\d{4}-\d{2}-\d{2}\b',
            r'\b\d{2}/\d{2}/\d{4}\b',
            r'\b(?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)[a-z]* \d{1,2},? \d{4}\b',
            r'\b\d{4}\b',
        ]
        
        for pattern in date_patterns:
            matches = re.findall(pattern, content)
            if matches:
                for match in matches:
                    try:
                        for fmt in ['%Y-%m-%d', '%d/%m/%Y', '%m/%d/%Y', '%B %d, %Y', '%b %d, %Y', '%Y']:
                            try:
                                return datetime.strptime(match, fmt)
                            except ValueError:
                                continue
                    except (ValueError, TypeError):
                        continue
        
        return None
    
    def _extract_additional_metadata(self, content: str, metadata: DocumentMetadata):
        email_matches = re.findall(self.patterns['email'][0], content)
        if email_matches:
            metadata.custom_fields['emails'] = ", ".join(list(set(email_matches))[:5])
        
        url_matches = re.findall(self.patterns['url'][0], content)
        if url_matches:
            metadata.custom_fields['urls'] = ", ".join(list(set(url_matches))[:5])
        
        isbn_matches = re.findall(self.patterns['isbn'][0], content)
        if isbn_matches:
            metadata.custom_fields['isbn'] = isbn_matches[0]
        
        doi_matches = re.findall(self.patterns['doi'][0], content, re.IGNORECASE)
        if doi_matches:
            metadata.custom_fields['doi'] = doi_matches[0]
    
    def _count_words_in_file(self, file_path: Path) -> int:
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
                words = re.findall(r'\b\w+\b', content)
                return len(words)
        except Exception:
            return 0
    
    def _detect_language(self, file_path: Path) -> str:
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read(5000)
                
                language_scores = {}
                for lang, patterns in self.language_patterns.items():
                    score = 0
                    for pattern in patterns:
                        matches = re.findall(pattern, content, re.IGNORECASE)
                        score += len(matches)
                    language_scores[lang] = score
                
                if language_scores:
                    max_lang = max(language_scores.items(), key=lambda x: x[1])
                    if max_lang[1] > 0:
                        return max_lang[0]
                
        except Exception:
            pass
        
        return 'en'
    
    def _extract_keywords(self, file_path: Path, title: Optional[str] = None) -> str:
        keywords = set()
        
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read(10000)
            
            if title:
                title_words = re.findall(r'\b[A-Za-z]{4,}\b', title)
                keywords.update(w.lower() for w in title_words)
            
            capitalized_phrases = re.findall(r'\b[A-Z][a-z]+(?:\s+[A-Z][a-z]+)+\b', content)
            keywords.update(phrase.lower() for phrase in capitalized_phrases[:10])
            
            words = re.findall(r'\b[a-z]{4,}\b', content.lower())
            word_freq = {}
            for word in words:
                word_freq[word] = word_freq.get(word, 0) + 1
            
            common_words = {'the', 'and', 'that', 'for', 'with', 'this', 'have', 'from', 'they', 'which',
                           'what', 'when', 'where', 'who', 'why', 'how', 'all', 'any', 'both', 'each',
                           'few', 'more', 'most', 'other', 'some', 'such', 'no', 'nor', 'not', 'only',
                           'own', 'same', 'so', 'than', 'too', 'very', 'just', 'now'}
            
            for word, freq in sorted(word_freq.items(), key=lambda x: x[1], reverse=True)[:30]:
                if freq >= 3 and word not in common_words:
                    keywords.add(word)
            
        except Exception:
            pass
        
        keyword_list = list(keywords)[:10]
        return ", ".join(keyword_list)


def extract_metadata(file_path: Union[str, Path], config: Optional[Dict] = None) -> DocumentMetadata:
    extractor = MetadataExtractor(config)
    return extractor.extract_metadata(file_path)


def batch_extract_metadata(file_paths: List[Union[str, Path]], config: Optional[Dict] = None) -> List[DocumentMetadata]:
    extractor = MetadataExtractor(config)
    results = []
    
    for file_path in file_paths:
        try:
            metadata = extractor.extract_metadata(file_path)
            results.append(metadata)
        except Exception as e:
            logger.error(f"Failed to extract metadata from {file_path}: {e}")
    
    return results


def metadata_to_json(metadata: DocumentMetadata, indent: int = 2) -> str:
    return metadata.to_json()


def metadata_from_json(json_str: str) -> DocumentMetadata:
    data = json.loads(json_str)
    return DocumentMetadata.from_dict(data)


def test_metadata_extraction() -> Dict[str, bool]:
    import tempfile
    
    test_results = {}
    
    with tempfile.NamedTemporaryFile(mode='w', suffix='.txt', delete=False, encoding='utf-8') as f:
        f.write("""Test Document Title

By John Doe

Created: 2023-10-15

This is a test document for metadata extraction.
It contains some sample content for testing purposes.

Keywords: testing, metadata, extraction, python""")
        txt_file = f.name
    
    try:
        metadata = extract_metadata(txt_file)
        test_results['txt_extraction'] = (
            metadata.content_title == "Test Document Title" and
            metadata.content_author == "John Doe" and
            metadata.core_file_type == "txt"
        )
    except Exception as e:
        test_results['txt_extraction'] = False
        logger.error(f"TXT extraction test failed: {e}")
    
    try:
        os.unlink(txt_file)
    except:
        pass
    
    return test_results


if __name__ == "__main__":
    import argparse
    
    parser = argparse.ArgumentParser(description="Metadata Extraction Module")
    parser.add_argument("--test", action="store_true", help="Run unit tests")
    parser.add_argument("--extract", type=str, help="Extract metadata from file")
    parser.add_argument("--batch", type=str, nargs='+', help="Extract metadata from multiple files")
    parser.add_argument("--output", type=str, choices=['json', 'yaml', 'text'], default='text',
                       help="Output format")
    parser.add_argument("--verbose", action="store_true", help="Verbose output")
    
    args = parser.parse_args()
    
    if args.verbose:
        logging.basicConfig(level=logging.INFO)
    
    if args.test:
        print("Running metadata extraction tests...")
        results = test_metadata_extraction()
        passed = sum(results.values())
        total = len(results)
        print(f"Tests passed: {passed}/{total} ({passed/total*100:.1f}%)")
        
        for test, success in results.items():
            status = "✓" if success else "✗"
            print(f"  {status} {test}")
    
    elif args.extract:
        file_path = Path(args.extract)
        
        if not file_path.exists():
            print(f"Error: File not found: {file_path}")
            exit(1)
        
        print(f"Extracting metadata from: {file_path}")
        
        try:
            metadata = extract_metadata(file_path)
            
            if args.output == 'json':
                print(metadata.to_json())
            elif args.output == 'yaml':
                if HAS_YAML:
                    import yaml
                    print(yaml.dump(metadata.to_dict(), allow_unicode=True, default_flow_style=False))
                else:
                    print("YAML output requires PyYAML. Install with: pip install pyyaml")
            else:
                print("\n=== METADATA EXTRACTION RESULTS ===")
                print(f"File: {metadata.core_file_name}")
                print(f"Type: {metadata.core_file_type}")
                print(f"Size: {metadata.core_file_size:,} bytes")
                
                if metadata.content_title:
                    print(f"Title: {metadata.content_title}")
                if metadata.content_author:
                    print(f"Author: {metadata.content_author}")
                if metadata.content_date:
                    print(f"Date: {metadata.content_date}")
                if metadata.content_language:
                    print(f"Language: {metadata.content_language}")
                if metadata.content_keywords:
                    print(f"Keywords: {metadata.content_keywords}")
                if metadata.processing_word_count:
                    print(f"Word Count: {metadata.processing_word_count:,}")
                if metadata.processing_page_count:
                    print(f"Page Count: {metadata.processing_page_count}")
                
                print(f"Processing Status: {metadata.processing_status}")
                if metadata.processing_error:
                    print(f"Error: {metadata.processing_error}")
        
        except Exception as e:
            print(f"Error extracting metadata: {e}")
    
    elif args.batch:
        file_paths = [Path(f) for f in args.batch]
        
        print(f"Extracting metadata from {len(file_paths)} files...")
        
        results = batch_extract_metadata(file_paths)
        
        print(f"\n=== BATCH EXTRACTION SUMMARY ===")
        print(f"Total files: {len(file_paths)}")
        print(f"Successfully extracted: {len(results)}")
        
        for metadata in results:
            print(f"\n{metadata.core_file_name}:")
            if metadata.content_title:
                print(f"  Title: {metadata.content_title}")
            print(f"  Status: {metadata.processing_status}")
    
    else:
        print("=" * 70)
        print("METADATA EXTRACTION MODULE")
        print("=" * 70)
        print("\nUsage examples:")
        print("  python metadata.py --extract document.pdf")
        print("  python metadata.py --extract document.pdf --output json")
        print("  python metadata.py --batch file1.txt file2.pdf")
        print("  python metadata.py --test")
        print("\nAvailable document types:")
        print("  • PDF (.pdf)")
        print("  • DOCX (.docx, .doc)")
        print("  • TXT (.txt, .text)")
        print("  • EPUB (.epub)")
        print("  • Markdown (.md, .markdown)")
        print("  • HTML (.html, .htm)")
        print("  • JSON (.json)")
        print("  • YAML (.yaml, .yml)")
        print("  • PowerPoint (.pptx, .ppt)")
        print("  • Excel (.xlsx, .xls)")
        print("  • CSV (.csv)")
        print("  • Images (.jpg, .png, .gif, .bmp)")